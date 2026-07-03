"""
core/loaders.py — with file hash tracking to skip re-indexing
"""

import os
import hashlib
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

HASH_FILE = "./faiss_index/.indexed_hashes"


def _file_hash(path):
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def is_already_indexed(path):
    if not os.path.exists(HASH_FILE):
        return False
    fh = _file_hash(path)
    with open(HASH_FILE) as f:
        return fh in f.read().splitlines()


def mark_as_indexed(path):
    os.makedirs(os.path.dirname(HASH_FILE), exist_ok=True)
    with open(HASH_FILE, "a") as f:
        f.write(_file_hash(path) + "\n")


def load_file_dynamically(file_path: str) -> list[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    loaders = {".pdf": _pdf, ".docx": _docx, ".pptx": _pptx, ".xlsx": _xlsx, ".csv": _csv}
    if ext not in loaders:
        raise ValueError("Unsupported: " + ext)
    docs = loaders[ext](file_path)
    if not docs:
        raise ValueError("No text extracted from " + os.path.basename(file_path))
    return docs


def get_chunks(docs, chunk_size=500, chunk_overlap=50):
    return RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    ).split_documents(docs)


def _pdf(path):
    from langchain_community.document_loaders import PyMuPDFLoader
    return PyMuPDFLoader(path).load()

def _docx(path):
    from langchain_community.document_loaders import Docx2txtLoader
    return Docx2txtLoader(path).load()

def _pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides):
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if texts:
            docs.append(Document(page_content="\n".join(texts), metadata={"source": path, "slide": i+1}))
    return docs

def _xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    docs = []
    for sheet in wb.sheetnames:
        rows = [" | ".join(str(c) for c in row if c is not None) for row in wb[sheet].iter_rows(values_only=True)]
        rows = [r for r in rows if r.strip()]
        if rows:
            docs.append(Document(page_content="\n".join(rows), metadata={"source": path, "sheet": sheet}))
    return docs

def _csv(path):
    import csv
    with open(path, newline="", encoding="utf-8-sig") as f:
        rows = [" | ".join(c.strip() for c in row if c.strip()) for row in csv.reader(f)]
    return [Document(page_content="\n".join(r for r in rows if r), metadata={"source": path})]
